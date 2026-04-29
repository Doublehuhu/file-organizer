"""文件预览服务 - 内容提取、缩略图生成、流媒体"""

import hashlib
import os
from pathlib import Path
import io
from app.config import settings
from app.utils.security import resolve_safe_path, validate_path_exists
from app.utils.file_types import get_preview_type


def get_preview_content(file_path: str, max_size: int | None = None) -> dict:
    """提取文件预览内容"""
    if max_size is None:
        max_size = settings.MAX_PREVIEW_TEXT_SIZE

    path = resolve_safe_path(file_path)
    validate_path_exists(path)

    if path.is_dir():
        return {"type": "directory", "content": f"目录: {path.name}", "encoding": "utf-8", "total_chars": 0, "truncated": False}

    preview_type = get_preview_type(file_path)

    if preview_type == "text":
        return _extract_text(path, max_size)
    elif preview_type == "document":
        return _extract_document(path, max_size)
    elif preview_type == "pdf":
        return _extract_pdf(path, max_size)
    elif preview_type == "image":
        return {"type": "image", "content": "", "encoding": "", "total_chars": 0, "truncated": False}
    elif preview_type in ("video", "audio"):
        return {"type": preview_type, "content": "", "encoding": "", "total_chars": 0, "truncated": False}
    else:
        return {"type": "unknown", "content": f"不支持预览此文件类型: {path.suffix}", "encoding": "utf-8", "total_chars": 0, "truncated": False}


def _extract_text(path: Path, max_size: int) -> dict:
    """提取文本文件内容"""
    import chardet
    try:
        file_size = path.stat().st_size
        truncated = file_size > max_size

        with open(path, "rb") as f:
            raw = f.read(max_size) if truncated else f.read()

        result = chardet.detect(raw)
        encoding = result.get("encoding", "utf-8") or "utf-8"

        try:
            content = raw.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            content = raw.decode("utf-8", errors="replace")

        return {
            "type": "text",
            "content": content,
            "encoding": encoding,
            "total_chars": len(content),
            "truncated": truncated,
        }
    except OSError as e:
        return {"type": "text", "content": f"读取失败: {e}", "encoding": "utf-8", "total_chars": 0, "truncated": False}


def _extract_document(path: Path, max_size: int) -> dict:
    """提取Office文档内容"""
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            content = "\n".join(p.text for p in doc.paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            slides = []
            for slide in prs.slides:
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                slides.append("\n".join(texts))
            content = "\n---\n".join(slides)
        elif ext in (".xlsx", ".xls", ".csv"):
            import openpyxl
            import csv
            if ext == ".csv":
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    reader = csv.reader(f)
                    rows = [",".join(row) for row in list(reader)[:100]]
                    content = "\n".join(rows)
            else:
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                ws = wb.active
                rows = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= 100:
                        break
                    rows.append("\t".join(str(cell) if cell is not None else "" for cell in row))
                content = "\n".join(rows)
                wb.close()
        else:
            return {"type": "document", "content": f"不支持预览: {ext}", "encoding": "utf-8", "total_chars": 0, "truncated": False}

        truncated = len(content) > max_size
        if truncated:
            content = content[:max_size]
        return {"type": "document", "content": content, "encoding": "utf-8", "total_chars": len(content), "truncated": truncated}
    except Exception as e:
        return {"type": "document", "content": f"解析失败: {e}", "encoding": "utf-8", "total_chars": 0, "truncated": False}


def _extract_pdf(path: Path, max_size: int) -> dict:
    """提取PDF内容"""
    try:
        import fitz
        doc = fitz.open(str(path))
        pages = []
        total_chars = 0
        for page in doc:
            text = page.get_text()
            pages.append(text)
            total_chars += len(text)
            if total_chars > max_size:
                break
        doc.close()
        content = "\n---\n".join(pages)
        truncated = total_chars > max_size
        if truncated:
            content = content[:max_size]
        return {"type": "pdf", "content": content, "encoding": "utf-8", "total_chars": len(content), "truncated": truncated}
    except Exception as e:
        return {"type": "pdf", "content": f"PDF解析失败: {e}", "encoding": "utf-8", "total_chars": 0, "truncated": False}


def get_thumbnail(file_path: str, size: int = 256) -> bytes | None:
    """生成缩略图"""
    path = resolve_safe_path(file_path)
    validate_path_exists(path)

    # 缓存键
    stat = path.stat()
    cache_key = hashlib.md5(
        f"{path.resolve()}{stat.st_size}{stat.st_mtime}".encode()
    ).hexdigest()
    cache_file = settings.CACHE_DIR / f"{cache_key}.webp"

    if cache_file.exists():
        return cache_file.read_bytes()

    preview_type = get_preview_type(file_path)

    try:
        from PIL import Image
        if preview_type == "image":
            img = Image.open(str(path))
            img.thumbnail((size, size), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="WEBP", quality=80)
            data = buf.getvalue()
            cache_file.write_bytes(data)
            return data
        elif preview_type == "pdf":
            import fitz
            doc = fitz.open(str(path))
            if len(doc) > 0:
                page = doc[0]
                pix = page.get_pixmap(dpi=72)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img.thumbnail((size, size), Image.LANCZOS)
                buf = io.BytesIO()
                img.save(buf, format="WEBP", quality=80)
                data = buf.getvalue()
                cache_file.write_bytes(data)
                doc.close()
                return data
            doc.close()
    except Exception:
        pass
    return None


def get_metadata(file_path: str) -> dict:
    """获取文件元数据"""
    path = resolve_safe_path(file_path)
    validate_path_exists(path)
    preview_type = get_preview_type(file_path)
    metadata = {}

    try:
        stat = path.stat()
        metadata["size"] = stat.st_size
        metadata["modified"] = stat.st_mtime

        if preview_type == "image":
            from PIL import Image
            with Image.open(str(path)) as img:
                metadata["dimensions"] = {"width": img.width, "height": img.height}
                metadata["format"] = img.format
        elif preview_type == "pdf":
            import fitz
            doc = fitz.open(str(path))
            metadata["pages"] = len(doc)
            if doc.metadata:
                metadata["title"] = doc.metadata.get("title", "")
                metadata["author"] = doc.metadata.get("author", "")
            doc.close()
        elif preview_type == "document":
            if path.suffix.lower() == ".docx":
                from docx import Document
                doc = Document(str(path))
                metadata["paragraphs"] = len(doc.paragraphs)
        elif preview_type in ("video", "audio"):
            metadata["duration"] = 0
    except Exception:
        pass

    return {"file_type": preview_type, "metadata": metadata}
